#!/usr/bin/env python3
"""
Unified lexical overlap report generator (HTML + PowerPoint formats).

This script consolidates:
  - build_lexical_overlap_report.py (HTML output)
  - build_lexical_overlap_pptx.py (PowerPoint output)

Generates comprehensive reports on lexical overlap in concept prompts
and correlations with concept-probe alignment.

Usage:
    python generate_lexical_report.py --format html
    python generate_lexical_report.py --format pptx
    python generate_lexical_report.py --format both

Outputs:
    results/lexical/lexical_overlap_investigation/LEXICAL_OVERLAP_REPORT.html
    results/lexical/lexical_overlap_investigation/LEXICAL_OVERLAP_REPORT.pptx

Env: llama2_env
Rachel C. Metzgar, Feb 2026
"""

import os
import sys
import csv
import json
import argparse
import base64
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from scipy import stats as sp_stats

from config import config

# ========================== CONFIG ========================== #

LEX_CSV = Path(str(config.RESULTS.lexical / "lexical_distinctiveness.csv"))
CONTRAST_JSON = Path(str(config.RESULTS.root / "probes" / "alignment" / "summaries" / "alignment_stats.json"))
STANDALONE_JSON = Path(str(config.RESULTS.root / "probes" / "standalone_alignment" / "summaries" / "standalone_alignment_stats.json"))
OUT_DIR = Path(str(config.RESULTS.lexical / "lexical_overlap_investigation"))
FIG_DIR = OUT_DIR / "figures"

OUT_DIR.mkdir(parents=True, exist_ok=True)
FIG_DIR.mkdir(exist_ok=True)

CAT_COLORS = {
    "Human vs AI (General)": "#888888",
    "Mental":     "#4C72B0",
    "Physical":   "#55A868",
    "Pragmatic":  "#C44E52",
    "Bio Ctrl":   "#8172B2",
    "Shapes":     "#CCB974",
    "SysPrompt":  "#64B5CD",
}


# ========================== DATA LOADING ========================== #

def load_data():
    """Load all data needed for report generation."""
    print("Loading data...")

    # Load lexical distinctiveness
    lex_data = []
    if LEX_CSV.exists():
        with open(LEX_CSV) as f:
            reader = csv.DictReader(f)
            for row in reader:
                lex_data.append({
                    "dim_id": int(row["dim_id"]),
                    "dim_name": row["dim_name"],
                    "category": row["category"],
                    "jaccard": float(row["jaccard"]),
                    "lex_distinct": float(row["lexical_distinctiveness"]),
                    "pct_human_entity": float(row["pct_human_entity_words"]),
                    "pct_ai_entity": float(row["pct_ai_entity_words"]),
                    "alignment": float(row["alignment_projection"]),
                })

    # Load alignment stats
    contrast_stats = {}
    if CONTRAST_JSON.exists():
        with open(CONTRAST_JSON) as f:
            contrast_stats = json.load(f)

    standalone_stats = {}
    if STANDALONE_JSON.exists():
        with open(STANDALONE_JSON) as f:
            standalone_stats = json.load(f)

    return lex_data, contrast_stats, standalone_stats


# ========================== FIGURE GENERATION ========================== #

def create_scatter_plot(lex_data):
    """Create scatter plot of lexical distinctiveness vs alignment."""
    fig, ax = plt.subplots(figsize=(10, 7))

    # Extract data
    categories = [d["category"] for d in lex_data]
    lex_distinct = [d["lex_distinct"] for d in lex_data]
    alignment = [abs(d["alignment"]) for d in lex_data]
    dim_names = [d["dim_name"] for d in lex_data]

    # Plot by category
    for cat in set(categories):
        cat_idx = [i for i, c in enumerate(categories) if c == cat]
        ax.scatter(
            [lex_distinct[i] for i in cat_idx],
            [alignment[i] for i in cat_idx],
            label=cat,
            color=CAT_COLORS.get(cat, "#999999"),
            s=100,
            alpha=0.7,
            edgecolors='black',
            linewidth=0.5
        )

    # Compute correlation
    rho, p = sp_stats.spearmanr(lex_distinct, alignment)

    # Labels
    ax.set_xlabel("Lexical Distinctiveness (1 - Jaccard)", fontsize=12, fontweight='bold')
    ax.set_ylabel("|Alignment with Conversational Probe|", fontsize=12, fontweight='bold')
    ax.set_title(f"Lexical Distinctiveness vs Alignment\nSpearman ρ = {rho:.3f}, p = {p:.4f}",
                 fontsize=14, fontweight='bold')
    ax.legend(loc='best', framealpha=0.9, fontsize=9)
    ax.grid(alpha=0.3, linestyle=':', linewidth=0.5)

    plt.tight_layout()
    fig.savefig(FIG_DIR / "scatter_lex_vs_alignment.png", dpi=200, bbox_inches='tight')

    # Convert to base64 for HTML embedding
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    b64 = base64.b64encode(buf.read()).decode("utf-8")

    return b64, rho, p


def create_category_bars(lex_data):
    """Create bar chart of average lexical distinctiveness by category."""
    # Aggregate by category
    cat_data = {}
    for d in lex_data:
        cat = d["category"]
        if cat not in cat_data:
            cat_data[cat] = []
        cat_data[cat].append(d["lex_distinct"])

    categories = sorted(cat_data.keys())
    means = [np.mean(cat_data[cat]) for cat in categories]

    fig, ax = plt.subplots(figsize=(10, 6))

    colors = [CAT_COLORS.get(cat, "#999999") for cat in categories]
    ax.bar(range(len(means)), means, color=colors, alpha=0.8, edgecolor='black', linewidth=0.8)

    ax.set_xlabel("Category", fontsize=12, fontweight='bold')
    ax.set_ylabel("Mean Lexical Distinctiveness", fontsize=12, fontweight='bold')
    ax.set_title("Lexical Distinctiveness by Category", fontsize=14, fontweight='bold')
    ax.set_xticks(range(len(categories)))
    ax.set_xticklabels(categories, rotation=45, ha='right', fontsize=10)
    ax.grid(axis='y', alpha=0.3, linestyle=':', linewidth=0.5)

    plt.tight_layout()
    fig.savefig(FIG_DIR / "category_lex_distinctiveness.png", dpi=200, bbox_inches='tight')

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    b64 = base64.b64encode(buf.read()).decode("utf-8")

    return b64


# ========================== HTML OUTPUT ========================== #

def generate_html_report(lex_data, scatter_b64, cat_bars_b64, rho, p):
    """Generate HTML report."""
    print("Generating HTML report...")

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Lexical Overlap Investigation Report</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; background: #f5f5f5; }}
        .container {{ max-width: 1200px; margin: 0 auto; background: white; padding: 40px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); }}
        h1 {{ color: #2c3e50; border-bottom: 3px solid #3498db; padding-bottom: 10px; }}
        h2 {{ color: #34495e; margin-top: 40px; }}
        .summary {{ background: #ecf0f1; padding: 20px; border-radius: 5px; margin: 20px 0; }}
        table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
        th, td {{ border: 1px solid #ddd; padding: 12px; text-align: left; }}
        th {{ background: #3498db; color: white; font-weight: bold; }}
        tr:nth-child(even) {{ background: #f2f2f2; }}
        figure {{ margin: 30px 0; text-align: center; }}
        img {{ max-width: 100%; height: auto; box-shadow: 0 2px 8px rgba(0,0,0,0.15); }}
        figcaption {{ margin-top: 10px; font-style: italic; color: #666; }}
        .highlight {{ background: #fff3cd; padding: 2px 5px; border-radius: 3px; }}
    </style>
</head>
<body>
    <div class="container">
        <h1>📊 Lexical Overlap Investigation Report</h1>
        <p><strong>Experiment 3: Concept-Probe Alignment Analysis</strong></p>
        <p>Generated: {OUT_DIR}</p>

        <div class="summary">
            <h2>Executive Summary</h2>
            <p>This report investigates whether <span class="highlight">lexical overlap</span> between human and AI concept prompts
            explains concept-probe alignment results.</p>
            <p><strong>Key Finding:</strong> Spearman ρ = {rho:.3f}, p = {p:.4f}</p>
            <p>{"<strong>Significant correlation detected.</strong>" if p < 0.05 else "<strong>No significant correlation.</strong>"}</p>
        </div>

        <h2>1. Lexical Distinctiveness vs Alignment</h2>
        <figure>
            <img src="data:image/png;base64,{scatter_b64}" alt="Scatter plot">
            <figcaption>Figure 1: Scatter plot showing relationship between lexical distinctiveness and alignment strength.</figcaption>
        </figure>

        <h2>2. Category-Level Analysis</h2>
        <figure>
            <img src="data:image/png;base64,{cat_bars_b64}" alt="Category bars">
            <figcaption>Figure 2: Average lexical distinctiveness by concept category.</figcaption>
        </figure>

        <h2>3. Dimension-Level Results</h2>
        <table>
            <thead>
                <tr>
                    <th>Dim ID</th>
                    <th>Dimension</th>
                    <th>Category</th>
                    <th>Jaccard</th>
                    <th>Lex Distinct</th>
                    <th>|Alignment|</th>
                </tr>
            </thead>
            <tbody>
"""

    for d in sorted(lex_data, key=lambda x: x["lex_distinct"], reverse=True):
        html += f"""
                <tr>
                    <td>{d["dim_id"]}</td>
                    <td>{d["dim_name"]}</td>
                    <td>{d["category"]}</td>
                    <td>{d["jaccard"]:.4f}</td>
                    <td>{d["lex_distinct"]:.4f}</td>
                    <td>{abs(d["alignment"]):.4f}</td>
                </tr>
"""

    html += """
            </tbody>
        </table>

        <h2>Conclusion</h2>
        <p>This analysis examines whether lexical overlap confounds our concept-probe alignment results.
        Dimensions with high lexical overlap might show inflated alignment simply due to shared vocabulary
        rather than true conceptual alignment.</p>

    </div>
</body>
</html>
"""

    out_path = OUT_DIR / "LEXICAL_OVERLAP_REPORT.html"
    with open(out_path, "w") as f:
        f.write(html)

    print(f"✅ HTML report saved: {out_path}")


# ========================== PPTX OUTPUT ========================== #

def generate_pptx_report(lex_data, rho, p):
    """Generate PowerPoint report."""
    print("Generating PowerPoint report...")

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("⚠️  python-pptx not installed. Skipping PowerPoint generation.")
        print("   Install with: pip install python-pptx")
        return

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Lexical Overlap Investigation"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x2c, 0x3e, 0x50)
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_para = title_frame.add_paragraph()
    subtitle_para.text = "Experiment 3: Concept-Probe Alignment"
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = RGBColor(0x7f, 0x8c, 0x8d)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Scatter plot slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = f"Lexical Distinctiveness vs Alignment (ρ = {rho:.3f}, p = {p:.4f})"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True

    scatter_path = FIG_DIR / "scatter_lex_vs_alignment.png"
    if scatter_path.exists():
        slide.shapes.add_picture(str(scatter_path), Inches(1), Inches(1.2), width=Inches(8))

    # Category bars slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Lexical Distinctiveness by Category"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True

    cat_path = FIG_DIR / "category_lex_distinctiveness.png"
    if cat_path.exists():
        slide.shapes.add_picture(str(cat_path), Inches(1), Inches(1.2), width=Inches(8))

    out_path = OUT_DIR / "LEXICAL_OVERLAP_REPORT.pptx"
    prs.save(str(out_path))

    print(f"✅ PowerPoint report saved: {out_path}")


# ========================== MAIN ========================== #

def main():
    parser = argparse.ArgumentParser(description="Generate lexical overlap report")
    parser.add_argument("--format", required=True, choices=["html", "pptx", "both"],
                        help="Output format: html, pptx, or both")
    args = parser.parse_args()

    print("=" * 80)
    print("LEXICAL OVERLAP REPORT GENERATION")
    print(f"Format: {args.format.upper()}")
    print("=" * 80)

    # Load data
    lex_data, contrast_stats, standalone_stats = load_data()

    if not lex_data:
        print("ERROR: No lexical distinctiveness data found!")
        print(f"Run lexical_distinctiveness.py first to generate {LEX_CSV}")
        return

    # Generate figures
    print("\nGenerating figures...")
    scatter_b64, rho, p = create_scatter_plot(lex_data)
    cat_bars_b64 = create_category_bars(lex_data)

    # Generate reports
    if args.format in ("html", "both"):
        generate_html_report(lex_data, scatter_b64, cat_bars_b64, rho, p)

    if args.format in ("pptx", "both"):
        generate_pptx_report(lex_data, rho, p)

    print("\n✅ Report generation complete!")


if __name__ == "__main__":
    main()
