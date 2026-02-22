#!/usr/bin/env python3
"""
Build a PowerPoint presentation on the Lexical Overlap Investigation.

Embeds figures from:
  - results/lexical_overlap_investigation/figures/  (newly generated)
  - results/concept_probe_alignment/figures/         (existing contrast)
  - results/standalone_alignment/figures/             (existing standalone)

Usage:
    python build_lexical_overlap_pptx.py

Output:
    results/lexical_overlap_investigation/LEXICAL_OVERLAP_REPORT.pptx
"""

import csv
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── paths ─────────────────────────────────────────────────────────────────
BASE = Path("/mnt/cup/labs/graziano/rachel/ai_mind_rep/exp_3/llama_exp_3-13B-chat")
LEX_FIG = BASE / "results" / "lexical_overlap_investigation" / "figures"
CONTRAST_FIG = BASE / "results" / "concept_probe_alignment" / "figures"
STANDALONE_FIG = BASE / "results" / "standalone_alignment" / "figures"
LEX_CSV = BASE / "results" / "concept_probe_alignment" / "summaries" / "lexical_distinctiveness.csv"
OUT_PATH = BASE / "results" / "lexical_overlap_investigation" / "LEXICAL_OVERLAP_REPORT.pptx"

# ── colors ────────────────────────────────────────────────────────────────
BLUE = RGBColor(0x4C, 0x72, 0xB0)
RED = RGBColor(0xC4, 0x4E, 0x52)
GREEN = RGBColor(0x55, 0xA8, 0x68)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARN_BG = RGBColor(0xFF, 0xF3, 0xE0)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)


# ── helpers ───────────────────────────────────────────────────────────────
def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLUE

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(12)

    return slide


def add_content_slide(prs, title, bullets=None, notes=""):
    """Add a slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # Bullets
    if bullets:
        body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(5.8))
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Support (text, kwargs) tuples for formatting
            if isinstance(bullet, tuple):
                text, kwargs = bullet
                p.text = text
                p.font.size = kwargs.get("size", Pt(16))
                if kwargs.get("bold"):
                    p.font.bold = True
                if kwargs.get("color"):
                    p.font.color.rgb = kwargs["color"]
                if kwargs.get("indent"):
                    p.level = kwargs["indent"]
                p.space_before = kwargs.get("space_before", Pt(6))
            else:
                p.text = bullet
                p.font.size = Pt(16)
                p.font.color.rgb = DARK
                p.space_before = Pt(6)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_figure_slide(prs, title, img_path, caption="", notes="",
                     img_left=None, img_top=None, img_width=None, img_height=None):
    """Add a slide with title and a figure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # Image
    if Path(img_path).exists():
        left = img_left or Inches(0.5)
        top = img_top or Inches(0.9)
        width = img_width or Inches(9.0)
        height = img_height or Inches(5.2)
        slide.shapes.add_picture(str(img_path), left, top, width, height)
    else:
        tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        tb.text_frame.paragraphs[0].text = f"[Figure not found: {img_path}]"
        tb.text_frame.paragraphs[0].font.color.rgb = RED

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9.0), Inches(0.6))
        tf = cap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.font.italic = True

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_two_figure_slide(prs, title, img_left_path, img_right_path,
                          caption_left="", caption_right="", notes=""):
    """Add a slide with two figures side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE

    for i, (path, cap) in enumerate([(img_left_path, caption_left),
                                      (img_right_path, caption_right)]):
        left = Inches(0.3 + i * 4.85)
        if Path(path).exists():
            slide.shapes.add_picture(str(path), left, Inches(1.0), Inches(4.5), Inches(4.5))
        if cap:
            cb = slide.shapes.add_textbox(left, Inches(5.6), Inches(4.5), Inches(0.8))
            tf = cb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cap
            p.font.size = Pt(9)
            p.font.color.rgb = GRAY
            p.font.italic = True

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_callout_slide(prs, title, heading, body, bg_color=LIGHT_BG):
    """Add a slide with a colored callout box."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # Callout box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.2), Inches(8.8), Inches(5.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.space_after = Pt(12)

    for line in body.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_before = Pt(6)

    return slide


# ═══════════════════════════════════════════════════════════════════════════
# Load data for table slides
# ═══════════════════════════════════════════════════════════════════════════
lex_rows = []
with open(LEX_CSV) as f:
    for row in csv.DictReader(f):
        lex_rows.append(row)


# ═══════════════════════════════════════════════════════════════════════════
# Build the presentation
# ═══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)


# ────────────────── SLIDE 1: Title ──────────────────
add_title_slide(
    prs,
    "Lexical Overlap Investigation",
    "Experiment 3: Concept-Probe Alignment in LLaMA-2-13B-Chat\n"
    "Could surface-level vocabulary overlap explain the alignment\n"
    "between concept directions and conversational probes?"
)


# ────────────────── SLIDE 2: Outline ──────────────────
add_content_slide(prs, "Outline", [
    ("1.  Background: The Concept-Probe Alignment Analysis", {"size": Pt(18), "bold": True}),
    ("2.  The Concern: Could Lexical Overlap Explain Our Results?", {"size": Pt(18), "bold": True}),
    ("3.  Six Lines of Evidence", {"size": Pt(18), "bold": True}),
    ("     (a) Vocabulary overlap between prompt conditions", {"size": Pt(15), "indent": 1, "color": GRAY}),
    ("     (b) Entity word contamination", {"size": Pt(15), "indent": 1, "color": GRAY}),
    ("     (c) Layer profiles", {"size": Pt(15), "indent": 1, "color": GRAY}),
    ("     (d) Standalone activations (no entity words)", {"size": Pt(15), "indent": 1, "color": GRAY}),
    ("     (e) Diagnostic dimensions (Formality, Shapes, SysPrompt)", {"size": Pt(15), "indent": 1, "color": GRAY}),
    ("     (f) Exp 2 conversation vocabulary", {"size": Pt(15), "indent": 1, "color": GRAY}),
    ("4.  Synthesis & Remaining Questions", {"size": Pt(18), "bold": True}),
])


# ────────────────── SLIDE 3: Background - Two Experiments ──────────────────
add_content_slide(prs, "Background: The Two Experiments", [
    ("Experiment 2: Conversational Partner-Identity Probes", {"size": Pt(18), "bold": True}),
    ("  A participant LLM (LLaMA-2-13B-Chat) had naturalistic conversations with", {"size": Pt(14)}),
    ("  either a human partner or an AI partner.", {"size": Pt(14)}),
    ("  Linear probes were trained on hidden states to classify: human vs. AI partner.", {"size": Pt(14)}),
    ("", {"size": Pt(6)}),
    ("Experiment 3: Concept-Probe Alignment", {"size": Pt(18), "bold": True}),
    ("  For 18 concept dimensions (emotions, agency, embodiment, shapes, etc.),", {"size": Pt(14)}),
    ("  we created 40 human-framed and 40 AI-framed prompts.", {"size": Pt(14)}),
    ("  Concept direction = mean(human activations) - mean(AI activations).", {"size": Pt(14)}),
    ("  We test whether each concept direction aligns with the Exp 2 probes.", {"size": Pt(14)}),
    ("", {"size": Pt(6)}),
    ("Key question: Does thinking about 'human emotions vs. AI emotions'", {"size": Pt(16), "bold": True, "color": BLUE}),
    ("activate the same neural direction that distinguishes partners in conversation?", {"size": Pt(16), "bold": True, "color": BLUE}),
])


# ────────────────── SLIDE 4: Probe Definitions ──────────────────
add_content_slide(prs, "How the Probes Work", [
    ("Control Probe", {"size": Pt(20), "bold": True, "color": BLUE}),
    ("  Extracts the hidden state at the last input token: [/INST]", {"size": Pt(14)}),
    ("  This is the boundary between the partner's last message and the", {"size": Pt(14)}),
    ("  participant's next response — the model is about to generate.", {"size": Pt(14)}),
    ("  Captures: active in-context representation of partner identity.", {"size": Pt(14), "bold": True}),
    ("", {"size": Pt(8)}),
    ("Reading Probe", {"size": Pt(20), "bold": True, "color": BLUE}),
    ('  Appends a reflective suffix: "I think the conversation partner of this user is"', {"size": Pt(14)}),
    ("  Placed in the participant's response position (after [/INST]).", {"size": Pt(14)}),
    ("  Extracts hidden state at the last token of this suffix.", {"size": Pt(14)}),
    ("  Captures: metacognitive reflection about partner identity.", {"size": Pt(14), "bold": True}),
    ("", {"size": Pt(8)}),
    ("Both probes: per-layer binary classifiers (Linear + Sigmoid, BCE loss)", {"size": Pt(13), "color": GRAY}),
    ("  41 layers, Human = 1, AI = 0. Trained on Exp 2 conversation data.", {"size": Pt(13), "color": GRAY}),
    ("Important: attention compresses the ENTIRE conversation into each position.", {"size": Pt(14), "bold": True, "color": RED}),
],
notes="Control probe: at [/INST] token, right before participant generates. "
      "Reading probe: at end of appended suffix 'I think the conversation partner of this user is'. "
      "Both extract from a single token position, but transformer attention means the entire "
      "conversation context is compressed into that representation.")


# ────────────────── SLIDE 5: Prompt Examples ──────────────────
add_content_slide(prs, "Example Prompts: Emotions Dimension", [
    ("Human-framed (label = 1):", {"size": Pt(16), "bold": True, "color": BLUE}),
    ('  "Imagine a human experiencing a sudden wave of fear when they', {"size": Pt(14)}),
    ('   hear an unexpected noise at night."', {"size": Pt(14)}),
    ("", {"size": Pt(6)}),
    ("AI-framed (label = 0):", {"size": Pt(16), "bold": True, "color": RED}),
    ('  "Imagine an AI system generating a high-threat classification when', {"size": Pt(14)}),
    ('   it detects an anomalous input signal."', {"size": Pt(14)}),
    ("", {"size": Pt(6)}),
    ("Standalone (no label — used in standalone analysis):", {"size": Pt(16), "bold": True, "color": GRAY}),
    ('  "Imagine experiencing a sudden wave of fear triggered by an', {"size": Pt(14)}),
    ('   unexpected noise at night."', {"size": Pt(14)}),
    ("", {"size": Pt(10)}),
    ("Concept direction = mean(human activations) - mean(AI activations)", {"size": Pt(15), "bold": True}),
    ("  Alignment = projection of concept direction onto Exp 2 probe weights", {"size": Pt(14)}),
])


# ────────────────── SLIDE 6: Main Result ──────────────────
add_figure_slide(
    prs, "Main Result: 15 of 18 Dimensions Align Significantly",
    CONTRAST_FIG / "fig_main_result.png",
    "Alignment between concept directions (Exp 3) and conversational partner probes (Exp 2). "
    "Tested via permutation (10K shuffles). Source: results/concept_probe_alignment/figures/fig_main_result.png",
    notes="15/18 dims significant at FDR-corrected p < 0.05. Mental and physical concepts project "
          "toward 'human' side. Shapes (negative control) and Formality show near-zero alignment."
)


# ────────────────── SLIDE 7: The Concern ──────────────────
add_callout_slide(
    prs,
    "The Concern: Lexical Overlap as a Confound",
    "Could surface-level vocabulary overlap explain the alignment?",
    "Type A: Entity Word Contamination\n"
    "  Human-framed prompts say 'human'; AI-framed say 'AI', 'machine', 'system'.\n"
    "  If the Exp 2 probe is just detecting these entity words, the alignment\n"
    "  would be trivially explained.\n"
    "\n"
    "Type B: Concept Vocabulary Overlap\n"
    "  Human emotion prompts use emotional vocabulary ('fear', 'joy').\n"
    "  If human-partner conversations also have more emotion words,\n"
    "  the probe might match word distributions, not conceptual content.\n"
    "\n"
    "We investigate with six independent lines of evidence.",
    bg_color=RGBColor(0xFF, 0xF0, 0xF0),
)


# ────────────────── SLIDE 8: Evidence 1 - Vocabulary Overlap ──────────────────
add_content_slide(prs, "Evidence 1: Vocabulary Overlap Between Prompt Conditions", [
    ("Method: Jaccard similarity between human-framed and AI-framed prompt vocabularies", {"size": Pt(15), "bold": True}),
    ("", {"size": Pt(4)}),
    ("  Jaccard = |intersection| / |union| of unique content words", {"size": Pt(14)}),
    ("  Lexical distinctiveness = 1 - Jaccard", {"size": Pt(14)}),
    ("  High distinctiveness = very different vocabulary between conditions", {"size": Pt(14)}),
    ("", {"size": Pt(8)}),
    ("Key question: Does vocabulary distinctiveness predict alignment?", {"size": Pt(16), "bold": True, "color": BLUE}),
    ("  If lexical overlap drives alignment, dimensions with more distinctive", {"size": Pt(14)}),
    ("  vocabulary should show stronger alignment.", {"size": Pt(14)}),
    ("", {"size": Pt(8)}),
    ("Data source: lexical_distinctiveness.csv", {"size": Pt(12), "color": GRAY}),
    ("Code: lexical_distinctiveness.py", {"size": Pt(12), "color": GRAY}),
])


# ────────────────── SLIDE 9: Jaccard bars ──────────────────
add_figure_slide(
    prs, "Vocabulary Overlap Is Low for Most Dimensions",
    LEX_FIG / "fig_jaccard_bars.png",
    "Jaccard similarity between human and AI prompt vocabularies. Most dimensions < 0.2. "
    "Baseline (dim 0) is 0.90 because prompts differ only in entity labels.",
)


# ────────────────── SLIDE 10: Lexical distinctiveness vs alignment ──────────────────
add_figure_slide(
    prs, "Lexical Distinctiveness Does NOT Predict Alignment",
    LEX_FIG / "fig_lexical_vs_alignment.png",
    "Spearman rho = 0.34, p = 0.17 (not significant). Vocabulary difference alone does not "
    "drive alignment. Data source: lexical_distinctiveness.csv",
)


# ────────────────── SLIDE 11: Evidence 2 heading ──────────────────
add_content_slide(prs, "Evidence 2: Entity Word Contamination", [
    ("Entity words = 'human', 'person', 'people' (human side)", {"size": Pt(15)}),
    ("                   'ai', 'artificial', 'machine', 'robot', 'system' (AI side)", {"size": Pt(15)}),
    ("", {"size": Pt(6)}),
    ("For each dimension, we computed:", {"size": Pt(15), "bold": True}),
    ("  % of human-framed prompts containing human entity words", {"size": Pt(14)}),
    ("  % of AI-framed prompts containing AI entity words", {"size": Pt(14)}),
    ("", {"size": Pt(8)}),
    ("Key test: Does entity contamination predict alignment magnitude?", {"size": Pt(16), "bold": True, "color": BLUE}),
])


# ────────────────── SLIDE 12: Entity contamination bars ──────────────────
add_figure_slide(
    prs, "Entity Contamination by Dimension",
    LEX_FIG / "fig_entity_contamination_bars.png",
    "Most Mental/Physical dims: 100% contamination on both sides. "
    "Key exceptions: Formality (25%/0%), Helpfulness (60%/0%), Shapes (0%/0%).",
)


# ────────────────── SLIDE 13: Entity vs alignment ──────────────────
add_figure_slide(
    prs, "Entity Contamination Correlates with Alignment — But This Is Expected",
    LEX_FIG / "fig_entity_vs_alignment.png",
    "Spearman rho = 0.66, p = 0.003. Significant, but entity words carry genuine conceptual "
    "content — 'a human experiencing fear' IS about a human. The question is whether alignment "
    "comes ONLY from entity words or also from deeper structure.",
)


# ────────────────── SLIDE 14: Evidence 3 - Layer Profiles ──────────────────
add_content_slide(prs, "Evidence 3: Layer Profiles Rule Out Surface Features", [
    ("If alignment were driven by surface-level vocabulary:", {"size": Pt(16), "bold": True}),
    ("  It should appear in EARLY layers (0-5), where representations", {"size": Pt(15)}),
    ("  are closest to raw token embeddings.", {"size": Pt(15)}),
    ("", {"size": Pt(8)}),
    ("What we observe:", {"size": Pt(16), "bold": True, "color": GREEN}),
    ("  Alignment peaks in LATE layers (25-40), where representations", {"size": Pt(15)}),
    ("  encode abstract, semantic content.", {"size": Pt(15)}),
    ("", {"size": Pt(8)}),
    ("This is inconsistent with a lexical/surface-feature explanation.", {"size": Pt(16), "bold": True}),
    ("", {"size": Pt(6)}),
    ("Note: even though the probe extracts from a single token position,", {"size": Pt(13), "color": GRAY}),
    ("the entire context is compressed into that representation via attention.", {"size": Pt(13), "color": GRAY}),
    ("The layer argument holds because early vs. late layers process differently,", {"size": Pt(13), "color": GRAY}),
    ("regardless of how information arrives at each position.", {"size": Pt(13), "color": GRAY}),
])


# ────────────────── SLIDE 15: Layer profiles figure ──────────────────
add_figure_slide(
    prs, "Layer Profiles: Selected Dimensions",
    LEX_FIG / "fig_layer_profiles_selected.png",
    "Per-layer alignment for Emotions (strong), Formality (near-zero), Shapes (near-zero). "
    "Alignment rises in late layers. Data: alignment_stats.json per-layer cosines.",
)


# ────────────────── SLIDE 16: Heatmap ──────────────────
add_figure_slide(
    prs, "Layer-wise Alignment Heatmap (Control Probe)",
    CONTRAST_FIG / "layerwise" / "fig_heatmap_control.png",
    "All 18 dimensions x 41 layers. Strong alignment (dark) concentrated in late layers (right side). "
    "Source: results/concept_probe_alignment/figures/layerwise/fig_heatmap_control.png",
)


# ────────────────── SLIDE 17: Layer profiles grid ──────────────────
add_figure_slide(
    prs, "Full Layer Profiles: All 18 Dimensions (Control Probe)",
    CONTRAST_FIG / "control_probe" / "fig_layer_profiles_grid.png",
    "Each panel = one dimension. Consistent pattern: flat/noisy early, rising late. "
    "Source: results/concept_probe_alignment/figures/control_probe/fig_layer_profiles_grid.png",
)


# ────────────────── SLIDE 18: Evidence 4 - Standalone ──────────────────
add_content_slide(prs, "Evidence 4: Standalone Activations (No Entity Words)", [
    ("Standalone prompts = same concepts, NO human/AI framing", {"size": Pt(16), "bold": True}),
    ("", {"size": Pt(6)}),
    ('  Contrast:    "Imagine a human experiencing a sudden wave of fear..."', {"size": Pt(14)}),
    ('  Standalone:  "Imagine experiencing a sudden wave of fear..."', {"size": Pt(14)}),
    ("", {"size": Pt(8)}),
    ("These prompts contain NO entity words at all.", {"size": Pt(15), "bold": True, "color": GREEN}),
    ("", {"size": Pt(6)}),
    ("We project standalone activations onto the Exp 2 probes:", {"size": Pt(15)}),
    ("  Does thinking about emotion IN GENERAL (no human/AI frame)", {"size": Pt(14)}),
    ("  activate a particular side of the human/AI probe?", {"size": Pt(14)}),
    ("", {"size": Pt(8)}),
    ("Important caveat: All standalone projections are negative", {"size": Pt(15), "bold": True, "color": RED}),
    ("  (baseline offset ~-1.5 to -1.96). The relative ranking is meaningful,", {"size": Pt(14)}),
    ("  not the absolute values.", {"size": Pt(14)}),
    ("", {"size": Pt(4)}),
    ("Statistical test: bootstrap against zero (10K resamples), BH-FDR corrected.", {"size": Pt(12), "color": GRAY}),
])


# ────────────────── SLIDE 19: Standalone ranked bars ──────────────────
add_figure_slide(
    prs, "Standalone Projections: Relative Ranking Matches Contrast",
    STANDALONE_FIG / "control_probe" / "fig_ranked_bars_all_layers.png",
    "All values negative (baseline offset). Mental/emotional concepts are least negative (toward human). "
    "Shapes most negative (toward AI). Source: results/standalone_alignment/figures/",
)


# ────────────────── SLIDE 20: Entity comparison ──────────────────
add_figure_slide(
    prs, "Standalone: Human Entity vs. AI Entity Concepts",
    STANDALONE_FIG / "standalone_specific" / "fig_entity_comparison.png",
    "'Think about what it means to be a human' (dim 16) vs. 'think about what it means to be an AI' (dim 17). "
    "Human entity projects closer to human side, as predicted.",
    img_width=Inches(8.0), img_height=Inches(5.0),
)


# ────────────────── SLIDE 21: Contrast vs standalone ──────────────────
add_figure_slide(
    prs, "Contrast vs. Standalone Alignment: Same Conceptual Structure?",
    LEX_FIG / "fig_contrast_vs_standalone.png",
    "Dimensions that align strongly in contrast analysis also show relatively less negative "
    "standalone projections. Data: alignment_stats.json + standalone_alignment_stats.json",
)


# ────────────────── SLIDE 22: Evidence 5 - Diagnostic dims ──────────────────
add_content_slide(prs, "Evidence 5: Diagnostic Dimensions", [
    ("Three dimensions serve as natural 'diagnostic tests':", {"size": Pt(16), "bold": True}),
    ("", {"size": Pt(6)}),
    ("Formality (dim 11): Different vocabulary, NO alignment", {"size": Pt(17), "bold": True, "color": BLUE}),
    ("  Human = casual language; AI = formal language", {"size": Pt(14)}),
    ("  Lexical distinctiveness = 0.867 (very different words)", {"size": Pt(14)}),
    ("  Alignment = -0.028 (n.s.) — vocabulary difference does NOT produce alignment", {"size": Pt(14), "bold": True}),
    ("", {"size": Pt(6)}),
    ("Shapes (dim 15): Negative control, ZERO alignment", {"size": Pt(17), "bold": True, "color": BLUE}),
    ("  Round pebbles vs angular objects — no human/AI relevance", {"size": Pt(14)}),
    ("  Entity contamination = 0% on both sides", {"size": Pt(14)}),
    ("  Alignment = 0.010 (n.s.) — confirms specificity to human/AI concepts", {"size": Pt(14), "bold": True}),
    ("", {"size": Pt(6)}),
    ("SysPrompt (dim 18): Pure entity labels, STRONG alignment", {"size": Pt(17), "bold": True, "color": BLUE}),
    ('  "You are talking to a human" vs "You are talking to an AI"', {"size": Pt(14)}),
    ("  Alignment = 0.215 (***) — entity words carry real conceptual weight", {"size": Pt(14), "bold": True}),
])


# ────────────────── SLIDE 23: SysPrompt variants ──────────────────
add_figure_slide(
    prs, "Standalone SysPrompt Variants: Human vs. AI References",
    STANDALONE_FIG / "standalone_specific" / "fig_sysprompt_variants.png",
    "4 SysPrompt variants: 'talk-to human' and 'bare human' project toward human side; "
    "'talk-to AI' and 'bare AI' toward AI side. Source: results/standalone_alignment/figures/",
    img_width=Inches(8.0), img_height=Inches(5.0),
)


# ────────────────── SLIDE 24: Evidence 6 - Exp 2 vocab ──────────────────
add_content_slide(prs, "Evidence 6: Exp 2 Conversation Vocabulary", [
    ("For Type B contamination, Exp 2 conversations would need asymmetric vocabulary", {"size": Pt(15), "bold": True}),
    ("  (e.g., more emotion words when talking to humans).", {"size": Pt(15)}),
    ("", {"size": Pt(6)}),
    ("We analyzed word frequencies across conditions in the participant's utterances:", {"size": Pt(15)}),
    ("", {"size": Pt(4)}),
    ("  Emotion words:          human-partner ~52/10K,  AI-partner ~47/10K,  ratio ~1.1", {"size": Pt(13)}),
    ("  Mental-state words:     human-partner ~89/10K,  AI-partner ~82/10K,  ratio ~1.1", {"size": Pt(13)}),
    ("  Embodiment words:       human-partner ~12/10K,  AI-partner ~11/10K,  ratio ~1.1", {"size": Pt(13)}),
    ("  Formal/technical words: human-partner ~18/10K,  AI-partner ~19/10K,  ratio ~0.95", {"size": Pt(13)}),
    ("  Identity labels:        human-partner ~5/10K,   AI-partner ~28/10K,  ratio ~0.18", {"size": Pt(13), "bold": True, "color": RED}),
    ("", {"size": Pt(8)}),
    ("Content vocabulary (emotions, mental states, etc.) is BALANCED across conditions.", {"size": Pt(16), "bold": True, "color": GREEN}),
    ("Only identity labels are strongly asymmetric (expected — AI convos mention 'AI' more).", {"size": Pt(14)}),
    ("The probe is NOT a word-frequency detector for emotion/mental vocabulary.", {"size": Pt(14), "bold": True}),
    ("", {"size": Pt(6)}),
    ("Code: vocab_asymmetry_check.py  |  Data: combined_all.csv (Exp 2)", {"size": Pt(12), "color": GRAY}),
])


# ────────────────── SLIDE 25: Synthesis ──────────────────
add_figure_slide(
    prs, "All Dimensions Ranked by Alignment",
    LEX_FIG / "fig_alignment_ranked.png",
    "The hierarchy (Mental > Physical > Pragmatic > Shapes/Formality) reflects conceptual "
    "relevance to human/AI identity, not vocabulary properties.",
)


# ────────────────── SLIDE 26: Summary table ──────────────────
add_content_slide(prs, "Summary: Six Lines of Evidence", [
    ("1. Lexical Distinctiveness", {"size": Pt(15), "bold": True}),
    ("   No correlation with alignment (rho = 0.34, p = 0.17)", {"size": Pt(13)}),
    ("2. Entity Contamination", {"size": Pt(15), "bold": True}),
    ("   Correlated (rho = 0.66, p = 0.003) but expected — entity words carry meaning", {"size": Pt(13)}),
    ("3. Layer Profiles", {"size": Pt(15), "bold": True}),
    ("   Alignment peaks in late layers (25-40), not early (lexical) layers", {"size": Pt(13)}),
    ("4. Standalone Activations", {"size": Pt(15), "bold": True}),
    ("   No entity words; relative ranking matches contrast analysis", {"size": Pt(13)}),
    ("5. Diagnostic Dimensions", {"size": Pt(15), "bold": True}),
    ("   Formality: diff vocab + no alignment. Shapes: no entity + no alignment.", {"size": Pt(13)}),
    ("6. Exp 2 Vocabulary", {"size": Pt(15), "bold": True}),
    ("   Concept-relevant words balanced across conditions (ratio ~1.1)", {"size": Pt(13)}),
])


# ────────────────── SLIDE 27: Conclusion ──────────────────
add_callout_slide(
    prs,
    "Conclusion",
    "Lexical overlap is a minor caveat, not a serious threat.",
    "Vocabulary distinctiveness does not predict alignment strength.\n"
    "Alignment concentrates in late model layers, not early (lexical) layers.\n"
    "Standalone prompts with no entity words still show differential patterns.\n"
    "Formality and Shapes (different vocab, no alignment) prove vocabulary\n"
    "  difference alone is not sufficient.\n"
    "Exp 2 conversation vocabulary is balanced for concept-relevant words.\n"
    "\n"
    "One legitimate caveat:\n"
    "  Entity word contamination correlates with alignment (rho = 0.66).\n"
    "  Some portion of the signal may come from entity-word matching,\n"
    "  but this is expected: entity words ARE part of the concept specification.",
    bg_color=GREEN_BG,
)


# ────────────────── SLIDE 28: Remaining Questions ──────────────────
add_content_slide(prs, "Remaining Questions", [
    ("1. Can we quantify the entity-word contribution?", {"size": Pt(16), "bold": True}),
    ("   A residual analysis (regress out entity contamination, test remaining", {"size": Pt(14)}),
    ("   alignment) could decompose the signal into entity vs. conceptual components.", {"size": Pt(14)}),
    ("", {"size": Pt(8)}),
    ("2. Does the standalone baseline offset have a lexical explanation?", {"size": Pt(16), "bold": True}),
    ("   All standalone prompts project to the AI side. Is this because short", {"size": Pt(14)}),
    ("   instructional text shares surface properties with AI-partner conversations?", {"size": Pt(14)}),
    ("", {"size": Pt(8)}),
    ("3. What drives the reading probe vs. control probe differences?", {"size": Pt(16), "bold": True}),
    ("   The reading probe uses a synthetic prompt never seen in training.", {"size": Pt(14)}),
    ("   Does it access different information than the control probe?", {"size": Pt(14)}),
    ("", {"size": Pt(8)}),
    ("4. What role does the system prompt play?", {"size": Pt(16), "bold": True}),
    ("   Exp 2 system prompts explicitly state partner identity.", {"size": Pt(14)}),
    ("   How much of the probe signal comes from system prompt vs. conversation?", {"size": Pt(14)}),
])


# ────────────────── SLIDE 29: Data Sources ──────────────────
add_content_slide(prs, "Data Sources and Code", [
    ("Scripts:", {"size": Pt(16), "bold": True}),
    ("  lexical_distinctiveness.py — Jaccard, entity contamination, correlations", {"size": Pt(12)}),
    ("  vocab_asymmetry_check.py — Exp 2 conversation vocabulary analysis", {"size": Pt(12)}),
    ("  2d_concept_probe_stats.py — contrast alignment (permutation tests)", {"size": Pt(12)}),
    ("  2e_concept_probe_figures.py — contrast figures", {"size": Pt(12)}),
    ("  3a_standalone_stats.py — standalone alignment (bootstrap tests)", {"size": Pt(12)}),
    ("  3b_standalone_figures.py — standalone figures", {"size": Pt(12)}),
    ("  build_lexical_overlap_report.py — generates HTML report + new figures", {"size": Pt(12)}),
    ("  build_lexical_overlap_pptx.py — generates this PowerPoint", {"size": Pt(12)}),
    ("", {"size": Pt(8)}),
    ("Data Files:", {"size": Pt(16), "bold": True}),
    ("  results/concept_probe_alignment/summaries/lexical_distinctiveness.csv", {"size": Pt(12)}),
    ("  results/concept_probe_alignment/summaries/alignment_stats.json", {"size": Pt(12)}),
    ("  results/standalone_alignment/summaries/standalone_alignment_stats.json", {"size": Pt(12)}),
    ("  data/concept_activations/contrasts/{dim}/concept_prompts.json", {"size": Pt(12)}),
    ("  data/concept_activations/standalone/{dim}/concept_prompts.json", {"size": Pt(12)}),
    ("", {"size": Pt(6)}),
    ("All paths relative to: llama_exp_3-13B-chat/", {"size": Pt(11), "color": GRAY}),
    ("Exp 2 probes trained on: llama_exp_2b-13B-chat/data/.../s###.csv (non-clean files)", {"size": Pt(11), "color": GRAY}),
])


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
prs.save(str(OUT_PATH))
size_mb = os.path.getsize(OUT_PATH) / (1024 * 1024)
print(f"Written: {OUT_PATH} ({size_mb:.1f} MB)")
print(f"  {len(prs.slides)} slides")
print("Done.")
